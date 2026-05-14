from __future__ import annotations

import json
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional

# Optional dependency:
# pip install python-pptx
from pptx import Presentation


logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(message)s"
)


@dataclass
class ShapeRecord:
    slide_index: int
    shape_key: str
    shape_name: str
    shape_type: str
    left: float
    top: float
    width: float
    height: float
    z_order: int
    group_path: List[str]
    has_text: bool
    text: str
    style_signature: Dict[str, Any]


def scan_pptx(pptx_path: str) -> Dict[str, Any]:
    prs = Presentation(pptx_path)
    inventory: Dict[str, Any] = {
        "source_pptx": pptx_path,
        "slides": []
    }

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_entry: Dict[str, Any] = {
            "slide_index": slide_index,
            "master_name": getattr(slide.slide_master, "name", None),
            "layout_name": getattr(slide.slide_layout, "name", None),
            "shapes": []
        }

        for shape_index, shape in enumerate(slide.shapes, start=1):
            has_text = bool(getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip())
            text_value = shape.text_frame.text if has_text else ""
            shape_entry: Dict[str, Any] = {
                "shape_key": f"{slide_index}/{getattr(shape, 'name', f'Shape{shape_index}')}",
                "shape_name": getattr(shape, "name", f"Shape{shape_index}"),
                "shape_type": str(getattr(shape, "shape_type", "unknown")),
                "left": float(getattr(shape, "left", 0)),
                "top": float(getattr(shape, "top", 0)),
                "width": float(getattr(shape, "width", 0)),
                "height": float(getattr(shape, "height", 0)),
                "z_order": shape_index,
                "group_path": [],
                "has_text": has_text,
                "text": text_value,
                "style_signature": {
                    "font_name": "KEEP_EXISTING",
                    "font_size": "KEEP_EXISTING",
                    "bold": "KEEP_EXISTING",
                    "italic": "KEEP_EXISTING",
                    "color": "KEEP_EXISTING",
                    "alignment": "KEEP_EXISTING"
                }
            }
            slide_entry["shapes"].append(shape_entry)

        inventory["slides"].append(slide_entry)

    return inventory


def load_json(path: str) -> Dict[str, Any]:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def save_json(obj: Dict[str, Any], path: str) -> None:
    Path(path).write_text(json.dumps(obj, indent=2, ensure_ascii=False), encoding="utf-8")


def replace_text_in_pptx(pptx_path: str, replacement_json_path: str, output_path: str) -> None:
    prs = Presentation(pptx_path)
    replacement_data = load_json(replacement_json_path)
    replacements = {item["shape_key"]: item["new_text"] for item in replacement_data.get("replacements", [])}

    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            shape_key = f"{slide_index}/{getattr(shape, 'name', '')}"
            if shape_key not in replacements:
                continue

            if not getattr(shape, "has_text_frame", False):
                continue

            # Important: keep formatting where possible.
            # This example updates the text content of the first paragraph/run path only.
            # A production implementation should preserve run-level formatting if mixed formatting exists.
            new_text = replacements[shape_key]
            text_frame = shape.text_frame

            # Safe baseline update
            text_frame.text = new_text

    prs.save(output_path)


def main():
    import argparse

    parser = argparse.ArgumentParser(description="PPTX scan and replace utility")
    parser.add_argument("--mode", choices=["scan", "replace"], required=True)
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--replacements", help="Replacement JSON file path")
    args = parser.parse_args()

    if args.mode == "scan":
        inventory = scan_pptx(args.pptx)
        save_json(inventory, args.output)
        logging.info("Inventory saved to %s", args.output)

    elif args.mode == "replace":
        if not args.replacements:
            raise ValueError("--replacements is required in replace mode")
        replace_text_in_pptx(args.pptx, args.replacements, args.output)
        logging.info("Updated PPTX saved to %s", args.output)


if __name__ == "__main__":
    main()
